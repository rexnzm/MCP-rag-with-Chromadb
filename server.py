"""
LangChain Chroma Vector DB Server with Ollama Embeddings and Multi-Format Document Ingestion.
Provides custom implementations for Ollama and OpenAI embeddings using direct API calls, 
and supports multiple document formats.
"""

import os
import uuid
import logging
import requests
import asyncio
from dotenv import load_dotenv, find_dotenv
from typing import List, Dict, Any, Tuple
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor
from langchain_chroma import Chroma
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from mcp.server.fastmcp import FastMCP

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from bs4 import BeautifulSoup
from pptx import Presentation
from openpyxl import load_workbook
from odf import text as odf_text, teletype
from odf.opendocument import load as odf_load
import markdown

# Load environment variables from .env file
_ = load_dotenv(find_dotenv())


# --- Logging ---
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("langchain_vector_db")


# --- Custom Ollama Embeddings (Direct API Implementation) ---
class DirectOllamaEmbeddings:
    """
    Custom Ollama embeddings implementation using direct API calls.
    This bypasses langchain_ollama connection issues in MCP context.
    Compatible with ChromaDB's embedding function interface.
    """

    def __init__(self, model: str, base_url: str = "http://localhost:11434"):
        """
        Initialize the embeddings with model name and Ollama base URL.

        Args:
            model: Name of the Ollama embedding model (e.g., "nomic-embed-text:latest")
            base_url: Base URL for Ollama API (default: http://localhost:11434)
        """
        self.model = model.replace(":latest", "")  # Ollama API prefers without :latest
        self.base_url = base_url.rstrip("/")
        self.api_url = f"{self.base_url}/api/embeddings"
        self._session = requests.Session()
        logger.info(f"Initialized DirectOllamaEmbeddings with model={self.model}, url={self.base_url}")

    def _get_embedding(self, text: str) -> List[float]:
        """
        Get embedding for a single text string.

        Args:
            text: Text to embed

        Returns:
            List of floats representing the embedding vector

        Raises:
            Exception: If the API call fails
        """
        try:
            payload = {
                "model": self.model,
                "prompt": text
            }

            response = self._session.post(
                self.api_url,
                json=payload,
                timeout=60
            )
            response.raise_for_status()

            data = response.json()
            embedding = data.get("embedding", [])

            if not embedding:
                raise ValueError(f"Empty embedding returned for text: {text[:100]}...")

            return embedding

        except requests.exceptions.ConnectionError as e:
            logger.error(f"Connection error to Ollama at {self.base_url}: {e}")
            raise Exception(f"Cannot connect to Ollama at {self.base_url}. Ensure Ollama is running.") from e
        except requests.exceptions.Timeout as e:
            logger.error(f"Timeout calling Ollama API: {e}")
            raise Exception("Ollama API request timed out. The model may be overloaded.") from e
        except requests.exceptions.HTTPError as e:
            logger.error(f"HTTP error from Ollama: {e}")
            raise Exception(f"Ollama API error: {e}") from e
        except Exception as e:
            logger.error(f"Error getting embedding: {e}")
            raise

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        """
        Embed a list of documents (required by LangChain interface).

        Args:
            texts: List of text strings to embed

        Returns:
            List of embedding vectors
        """
        embeddings = []
        for i, text in enumerate(texts):
            try:
                if not text or not text.strip():
                    logger.warning(f"Skipping empty text at index {i}")
                    # Return zero vector for empty texts
                    embeddings.append([0.0] * 768)  # nomic-embed-text dimension
                    continue

                embedding = self._get_embedding(text)
                embeddings.append(embedding)

                # Log progress for large batches
                if (i + 1) % 10 == 0:
                    logger.debug(f"Embedded {i + 1}/{len(texts)} documents")

            except Exception as e:
                logger.error(f"Failed to embed document {i}: {e}")
                raise

        logger.debug(f"Successfully embedded {len(embeddings)} documents")
        return embeddings

    def embed_query(self, text: str) -> List[float]:
        """
        Embed a single query text (required by LangChain interface).

        Args:
            text: Query text to embed

        Returns:
            Embedding vector
        """
        if not text or not text.strip():
            logger.warning("Embedding empty query text")
            return [0.0] * 768  # nomic-embed-text dimension

        return self._get_embedding(text)

    def __call__(self, text: str) -> List[float]:
        """
        Allow the class to be called directly (ChromaDB compatibility).

        Args:
            text: Text to embed

        Returns:
            Embedding vector
        """
        return self.embed_query(text)


# --- Custom OpenAI Embeddings (Direct API Implementation) ---
class DirectOpenAIEmbeddings:
    """
    Custom OpenAI embeddings implementation using direct API calls.
    Uses OpenAI's text-embedding-3-small model by default.
    Compatible with ChromaDB's embedding function interface.
    """

    def __init__(self, api_key: str, model: str = "text-embedding-3-small"):
        """
        Initialize the embeddings with OpenAI API key and model name.

        Args:
            api_key: OpenAI API key
            model: Name of the OpenAI embedding model (default: text-embedding-3-small)
                  Options: text-embedding-3-small (1536 dim), text-embedding-3-large (3072 dim)
        """
        if not api_key:
            raise ValueError("OpenAI API key is required. Set OPENAI_API_KEY environment variable.")

        self.api_key = api_key
        self.model = model
        self.api_url = "https://api.openai.com/v1/embeddings"
        self._session = requests.Session()
        self._session.headers.update({
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        })

        # Set embedding dimensions based on model
        self.dimensions = 1536 if "small" in model else 3072

        logger.info(f"Initialized DirectOpenAIEmbeddings with model={self.model}, dimensions={self.dimensions}")

    def _get_embedding(self, text: str) -> List[float]:
        """
        Get embedding for a single text string.

        Args:
            text: Text to embed

        Returns:
            List of floats representing the embedding vector

        Raises:
            Exception: If the API call fails
        """
        try:
            payload = {
                "model": self.model,
                "input": text,
                "encoding_format": "float"
            }

            response = self._session.post(
                self.api_url,
                json=payload,
                timeout=60
            )
            response.raise_for_status()

            data = response.json()

            # OpenAI returns embeddings in data[0].embedding
            if "data" not in data or len(data["data"]) == 0:
                raise ValueError(f"Empty embedding returned for text: {text[:100]}...")

            embedding = data["data"][0]["embedding"]

            if not embedding:
                raise ValueError(f"Empty embedding vector for text: {text[:100]}...")

            return embedding

        except requests.exceptions.ConnectionError as e:
            logger.error(f"Connection error to OpenAI API: {e}")
            raise Exception("Cannot connect to OpenAI API. Check your internet connection.") from e
        except requests.exceptions.Timeout as e:
            logger.error(f"Timeout calling OpenAI API: {e}")
            raise Exception("OpenAI API request timed out. Try again later.") from e
        except requests.exceptions.HTTPError as e:
            status_code = e.response.status_code if e.response else "unknown"
            error_msg = e.response.text if e.response else str(e)
            logger.error(f"HTTP error from OpenAI (status {status_code}): {error_msg}")

            if status_code == 401:
                raise Exception("Invalid OpenAI API key. Check your OPENAI_API_KEY environment variable.") from e
            elif status_code == 429:
                raise Exception("OpenAI API rate limit exceeded. Wait a moment and try again.") from e
            else:
                raise Exception(f"OpenAI API error (status {status_code}): {error_msg}") from e
        except Exception as e:
            logger.error(f"Error getting embedding: {e}")
            raise

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        """
        Embed a list of documents (required by LangChain interface).
        Uses batch API calls for efficiency when possible.

        Args:
            texts: List of text strings to embed

        Returns:
            List of embedding vectors
        """
        if not texts:
            return []

        embeddings = []

        # OpenAI supports batch embeddings, process in batches of 100
        batch_size = 100

        for batch_start in range(0, len(texts), batch_size):
            batch_end = min(batch_start + batch_size, len(texts))
            batch_texts = texts[batch_start:batch_end]

            try:
                # For batch processing, send all texts at once
                payload = {
                    "model": self.model,
                    "input": batch_texts,
                    "encoding_format": "float"
                }

                response = self._session.post(
                    self.api_url,
                    json=payload,
                    timeout=120  # Longer timeout for batches
                )
                response.raise_for_status()

                data = response.json()

                if "data" not in data:
                    raise ValueError("Invalid response from OpenAI API")

                # Extract embeddings in correct order
                batch_embeddings = [item["embedding"] for item in sorted(data["data"], key=lambda x: x["index"])]
                embeddings.extend(batch_embeddings)

                # Log progress for large batches
                if batch_end % 100 == 0 or batch_end == len(texts):
                    logger.debug(f"Embedded {batch_end}/{len(texts)} documents")

            except Exception as e:
                logger.error(f"Failed to embed batch {batch_start}-{batch_end}: {e}")
                raise

        logger.debug(f"Successfully embedded {len(embeddings)} documents")
        return embeddings

    def embed_query(self, text: str) -> List[float]:
        """
        Embed a single query text (required by LangChain interface).

        Args:
            text: Query text to embed

        Returns:
            Embedding vector
        """
        if not text or not text.strip():
            logger.warning("Embedding empty query text")
            return [0.0] * self.dimensions

        return self._get_embedding(text)

    def __call__(self, text: str) -> List[float]:
        """
        Allow the class to be called directly (ChromaDB compatibility).

        Args:
            text: Text to embed

        Returns:
            Embedding vector
        """
        return self.embed_query(text)


# --- Configuration ---
current_dir = Path(__file__).parent
CHROMA_PATH = str(current_dir.joinpath("chroma_db"))
CHUNK_SIZE = 4096
CHUNK_OVERLAP = CHUNK_SIZE // 10  # 10% overlap
COLLECTION_NAME = "documents"
DOWNLOAD_DIR = str(current_dir.joinpath("downloads"))
MAX_WORKERS = min(8, (os.cpu_count() or 1) * 2)

# Embedding provider selection (choose one: 'openai' or 'ollama')
EMBEDDING_PROVIDER = os.getenv("EMBEDDING_PROVIDER", "openai").lower()

# OpenAI Configuration
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
OPENAI_EMBED_MODEL = os.getenv("OPENAI_EMBED_MODEL", "text-embedding-3-small")

# Ollama Configuration (fallback option)
OLLAMA_EMBED_MODEL = os.getenv("OLLAMA_EMBED_MODEL", "nomic-embed-text")
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")

# Determine current embed model name used for diagnostics (fixes undefined name issues)
EMBED_MODEL = OLLAMA_EMBED_MODEL if EMBEDDING_PROVIDER == "ollama" else OPENAI_EMBED_MODEL

# Supported file extensions
SUPPORTED_EXTENSIONS = {
    '.pdf', '.docx', '.doc', '.txt', '.html', '.htm',
    '.pptx', '.ppt', '.xlsx', '.xls', '.odt', '.md',
    '.rtf', '.csv', '.json', '.xml'
}

# --- MCP ---
mcp = FastMCP("langchain-vector-db")

# --- Reusable resources ---
_session = requests.Session()

# Initialize embeddings based on provider selection
if EMBEDDING_PROVIDER == "openai":
    logger.info("Using OpenAI embeddings")
    if not OPENAI_API_KEY:
        logger.error("=" * 60)
        logger.error("OPENAI_API_KEY not set!")
        logger.error("Please set your OpenAI API key in .env file:")
        logger.error("OPENAI_API_KEY=your-api-key-here")
        logger.error("=" * 60)
        raise ValueError("OPENAI_API_KEY environment variable is required when using OpenAI embeddings")

    embeddings = DirectOpenAIEmbeddings(
        api_key=OPENAI_API_KEY,
        model=OPENAI_EMBED_MODEL
    )
elif EMBEDDING_PROVIDER == "ollama":
    logger.info("Using Ollama embeddings")
    embeddings = DirectOllamaEmbeddings(
        model=OLLAMA_EMBED_MODEL,
        base_url=OLLAMA_BASE_URL
    )
else:
    raise ValueError(f"Invalid EMBEDDING_PROVIDER: {EMBEDDING_PROVIDER}. Must be 'openai' or 'ollama'")

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP,
    length_function=len,
    separators=["\n\n", "\n", " ", ""]
)

vectorstore = Chroma(
    persist_directory=CHROMA_PATH,
    embedding_function=embeddings,
    collection_name=COLLECTION_NAME
)


# --- Helper Functions ---
def _short_uuid() -> str:
    """Generate a short UUID for use in chunk IDs."""
    return uuid.uuid4().hex[:8]


def check_ollama_health() -> Dict[str, Any]:
    """
    Check if Ollama service is running and the embedding model is available.
    Returns status information and suggestions for fixing issues.
    """
    health_status = {
        "service_running": False,
        "model_available": False,
        "can_generate_embeddings": False,
        "errors": [],
        "suggestions": []
    }

    # Check if service is running
    try:
        response = _session.get(f"{OLLAMA_BASE_URL}/api/tags", timeout=5)
        if response.status_code == 200:
            health_status["service_running"] = True
            data = response.json()
            models = data.get("models", [])
            model_names = [m.get("name", "") for m in models]

            # Check if our model is installed
            model_base = EMBED_MODEL.replace(":latest", "")
            if EMBED_MODEL in model_names or model_base in model_names or f"{model_base}:latest" in model_names:
                health_status["model_available"] = True
            else:
                health_status["errors"].append(f"Model '{EMBED_MODEL}' not found")
                health_status["suggestions"].append(f"Run: ollama pull {model_base}")
        else:
            health_status["errors"].append(f"Ollama responded with status {response.status_code}")
    except requests.exceptions.ConnectionError:
        health_status["errors"].append("Cannot connect to Ollama service")
        health_status["suggestions"].append("Make sure Ollama is running at " + OLLAMA_BASE_URL)
        health_status["suggestions"].append("Start Ollama service or run 'ollama serve'")
    except requests.exceptions.Timeout:
        health_status["errors"].append("Connection to Ollama timed out")
        health_status["suggestions"].append("Ollama may be overloaded - try restarting it")
    except Exception as e:
        health_status["errors"].append(f"Error connecting to Ollama: {str(e)}")

    # If service and model are OK, test embedding generation
    if health_status["service_running"] and health_status["model_available"]:
        try:
            payload = {
                "model": EMBED_MODEL.replace(":latest", ""),
                "prompt": "test"
            }
            response = _session.post(
                f"{OLLAMA_BASE_URL}/api/embeddings",
                json=payload,
                timeout=10
            )
            if response.status_code == 200 and response.json().get("embedding"):
                health_status["can_generate_embeddings"] = True
            else:
                health_status["errors"].append("Cannot generate embeddings")
                health_status["suggestions"].append("Try restarting Ollama service")
        except requests.exceptions.ConnectionError:
            health_status["errors"].append("Connection closed during embedding test")
            health_status["suggestions"].append("Ollama may have crashed - check system resources (RAM/CPU)")
            health_status["suggestions"].append("Try restarting Ollama or your computer")
        except Exception as e:
            health_status["errors"].append(f"Embedding test failed: {str(e)}")

    return health_status


def validate_ollama_connection() -> bool:
    """
    Validate Ollama connection and log detailed error messages if there are issues.
    Returns True if connection is healthy, False otherwise.
    """
    logger.info("Validating Ollama connection...")
    health = check_ollama_health()

    if health["can_generate_embeddings"]:
        logger.info("✓ Ollama connection validated successfully")
        return True
    else:
        logger.error("=" * 60)
        logger.error("OLLAMA CONNECTION ISSUES DETECTED")
        logger.error("=" * 60)

        for error in health["errors"]:
            logger.error(f"✗ {error}")

        if health["suggestions"]:
            logger.error("\nSuggested solutions:")
            for suggestion in health["suggestions"]:
                logger.error(f"  → {suggestion}")

        logger.error("\nFor detailed diagnostics, run:")
        logger.error("  python check_ollama.py")
        logger.error("=" * 60)

        return False


def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from PDF using PyPDF2. Returns empty string on failure."""
   
    try:
        reader = PdfReader(pdf_path)
        pages = []
        for page in reader.pages:
            try:
                pt = page.extract_text()
            except Exception:
                pt = None
            if pt:
                pages.append(pt)
        return "\n".join(pages)
    except Exception as exc:
        logger.exception("Failed to read PDF %s: %s", pdf_path, exc)
        return ""


def extract_text_from_docx(docx_path: str) -> str:
    """Extract text from DOCX files. Returns empty string on failure."""
   
    try:
        doc = DocxDocument(docx_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

        # Also extract text from tables
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    tables_text.append(" | ".join(row_text))

        all_text = paragraphs + tables_text
        return "\n".join(all_text)
    except Exception as exc:
        logger.exception("Failed to read DOCX %s: %s", docx_path, exc)
        return ""


def extract_text_from_txt(txt_path: str) -> str:
    """Extract text from TXT files. Returns empty string on failure."""
    try:
        with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as exc:
        logger.exception("Failed to read TXT %s: %s", txt_path, exc)
        return ""


def extract_text_from_html(html_path: str) -> str:
    """Extract text from HTML files. Returns empty string on failure."""
    
    try:
        with open(html_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
        soup = BeautifulSoup(html_content, 'html.parser')

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        text = soup.get_text(separator='\n', strip=True)
        # Clean up excessive newlines
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)
    except Exception as exc:
        logger.exception("Failed to read HTML %s: %s", html_path, exc)
        return ""


def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract text from PPTX files. Returns empty string on failure."""
    
    try:
        prs = Presentation(pptx_path)
        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_runs.append(shape.text)

        return "\n".join(text_runs)
    except Exception as exc:
        logger.exception("Failed to read PPTX %s: %s", pptx_path, exc)
        return ""


def extract_text_from_xlsx(xlsx_path: str) -> str:
    """Extract text from XLSX files. Returns empty string on failure."""
    
    try:
        workbook = load_workbook(xlsx_path, read_only=True, data_only=True)
        text_parts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===")

            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:
                    text_parts.append(" | ".join(row_text))

        workbook.close()
        return "\n".join(text_parts)
    except Exception as exc:
        logger.exception("Failed to read XLSX %s: %s", xlsx_path, exc)
        return ""


def extract_text_from_odt(odt_path: str) -> str:
    """Extract text from ODT files. Returns empty string on failure."""
    
    try:
        doc = odf_load(odt_path)
        all_paragraphs = doc.getElementsByType(odf_text.P)
        text_parts = [teletype.extractText(para) for para in all_paragraphs]
        return "\n".join([text for text in text_parts if text.strip()])
    except Exception as exc:
        logger.exception("Failed to read ODT %s: %s", odt_path, exc)
        return ""


def extract_text_from_markdown(md_path: str) -> str:
    """Extract text from Markdown files. Returns empty string on failure."""
    try:
        with open(md_path, 'r', encoding='utf-8', errors='ignore') as f:
            md_content = f.read()
            html = markdown.markdown(md_content)        # convert to plain text
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text(separator='\n', strip=True)

    except Exception as exc:
        logger.exception("Failed to read Markdown %s: %s", md_path, exc)
        return ""


def extract_text_from_csv(csv_path: str) -> str:
    """Extract text from CSV files. Returns empty string on failure."""
    try:
        import csv
        text_parts = []
        with open(csv_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for row in reader:
                if row:
                    text_parts.append(" | ".join(row))
        return "\n".join(text_parts)
    except Exception as exc:
        logger.exception("Failed to read CSV %s: %s", csv_path, exc)
        return ""


def extract_text_from_json(json_path: str) -> str:
    """Extract text from JSON files. Returns empty string on failure."""
    try:
        import json
        with open(json_path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
        # Convert JSON to formatted string
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as exc:
        logger.exception("Failed to read JSON %s: %s", json_path, exc)
        return ""


def extract_text_from_xml(xml_path: str) -> str:
    """Extract text from XML files. Returns empty string on failure."""
    
    try:
        with open(xml_path, 'r', encoding='utf-8', errors='ignore') as f:
            xml_content = f.read()
        soup = BeautifulSoup(xml_content, 'xml')
        return soup.get_text(separator='\n', strip=True)
    except Exception as exc:
        logger.exception("Failed to read XML %s: %s", xml_path, exc)
        return ""


def extract_text(file_path: str) -> str:
    """
    Extract text from a file based on its extension.
    Supports PDF, DOCX, TXT, HTML, PPTX, XLSX, ODT, MD, CSV, JSON, XML, and more.
    Returns empty string on failure.
    """
    file_ext = Path(file_path).suffix.lower()

    extractors = {
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.doc': extract_text_from_docx,  # Attempt DOCX extractor for .doc
        '.txt': extract_text_from_txt,
        '.text': extract_text_from_txt,
        '.html': extract_text_from_html,
        '.htm': extract_text_from_html,
        '.pptx': extract_text_from_pptx,
        '.ppt': extract_text_from_pptx,  # Attempt PPTX extractor for .ppt
        '.xlsx': extract_text_from_xlsx,
        '.xls': extract_text_from_xlsx,  # Attempt XLSX extractor for .xls
        '.odt': extract_text_from_odt,
        '.md': extract_text_from_markdown,
        '.markdown': extract_text_from_markdown,
        '.csv': extract_text_from_csv,
        '.json': extract_text_from_json,
        '.xml': extract_text_from_xml,
        '.rtf': extract_text_from_txt,  # Fallback to TXT for RTF
    }

    extractor = extractors.get(file_ext)
    if extractor:
        return extractor(file_path)
    else:
        logger.warning("Unsupported file type: %s, attempting TXT extraction", file_ext)
        return extract_text_from_txt(file_path)


def document_to_chunks(file_path: str) -> Tuple[List[Document], List[str]]:
    """
    Convert a document (any supported format) into chunk Documents and generated ids.
    Supports PDF, DOCX, TXT, HTML, PPTX, XLSX, ODT, MD, CSV, JSON, XML, and more.
    Does NOT write to vectorstore.
    """
    text = extract_text(file_path)
    if not text:
        logger.warning("No text extracted from %s", file_path)
        return [], []

    src_path = Path(file_path)
    doc = Document(
        page_content=text,
        metadata={
            "source": str(src_path),
            "filename": src_path.name,
            "file_type": src_path.suffix.lower()
        }
    )

    chunks = text_splitter.split_documents([doc])
    total = len(chunks)

    ids: List[str] = []
    for i, chunk in enumerate(chunks):
        chunk.metadata.setdefault("source", str(src_path))
        chunk.metadata.setdefault("filename", src_path.name)
        chunk.metadata.setdefault("file_type", src_path.suffix.lower())
        chunk.metadata["chunk_index"] = i
        chunk.metadata["total_chunks"] = total
        # include short uuid to avoid id collisions
        short = _short_uuid()
        ids.append(f"{src_path.stem}_chunk_{i}_{short}")

    return chunks, ids


# Legacy function for backwards compatibility
def pdf_to_chunks(pdf_path: str) -> Tuple[List[Document], List[str]]:
    """
    Legacy function. Use document_to_chunks instead.
    Convert a single PDF into chunk Documents and generated ids.
    """
    return document_to_chunks(pdf_path)


def download_file(url: str, download_dir: str = DOWNLOAD_DIR) -> str:
    """Download file from URL using a persistent session. Returns local path."""
    os.makedirs(download_dir, exist_ok=True)
    parsed_name = Path(url.split("?", 1)[0]).name

    # If no extension detected, try to infer from content-type later
    if not Path(parsed_name).suffix:
        parsed_name = f"downloaded_{uuid.uuid4().hex[:8]}"

    local_path = str(Path(download_dir).joinpath(parsed_name))

    try:
        with _session.get(url, stream=True, timeout=30) as resp:
            resp.raise_for_status()
            ct = resp.headers.get("content-type", "")
            logger.debug("download_file: content-type for %s is %s", url, ct)

            with open(local_path, "wb") as f:
                for chunk in resp.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)
    except Exception as exc:
        logger.exception("Failed to download %s: %s", url, exc)
        raise

    return local_path


# Legacy function for backwards compatibility
def download_pdf(url: str, download_dir: str = DOWNLOAD_DIR) -> str:
    """Legacy function. Use download_file instead. Download PDF from URL."""
    return download_file(url, download_dir)


# --- MCP Tools ---
@mcp.tool()
async def ingest_document(source: str) -> Dict[str, Any]:
    """
    Ingest document(s) from URL, folder, or file path.
    Supports multiple formats: PDF, DOCX, TXT, HTML, PPT/PPTX, XLSX, ODT, MD, CSV, JSON, XML, and more.
    Uses thread pool to process documents concurrently, batches DB writes for performance.
    """
    loop = asyncio.get_running_loop()
    total_chunks = 0
    processed_files: List[str] = []
    all_chunks: List[Document] = []
    all_ids: List[str] = []

    try:
        # URL case
        if source.startswith(("http://", "https://")):
            logger.info("Downloading file from URL: %s", source)
            try:
                local = await loop.run_in_executor(None, download_file, source)
            except Exception as exc:
                return {"status": "error", "message": f"Download failed: {exc}"}
            processed_files.append(local)
            chunks, ids = await loop.run_in_executor(None, document_to_chunks, local)
            all_chunks.extend(chunks)
            all_ids.extend(ids)
            total_chunks += len(chunks)

        # Directory case
        elif os.path.isdir(source):
            # Find all supported document files
            doc_paths = []
            for ext in SUPPORTED_EXTENSIONS:
                doc_paths.extend(Path(source).glob(f"*{ext}"))

            doc_paths = sorted([str(p) for p in doc_paths])

            if not doc_paths:
                return {"status": "error", "message": f"No supported document files found in folder: {source}"}
            logger.info("Found %d document files in %s", len(doc_paths), source)

            # Process documents concurrently but batch writes later
            with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
                tasks = [
                    loop.run_in_executor(executor, document_to_chunks, doc_path)
                    for doc_path in doc_paths
                ]
                for doc_path, fut in zip(doc_paths, await asyncio.gather(*tasks, return_exceptions=False)):
                    chunks, ids = fut  # type: ignore
                    if chunks:
                        all_chunks.extend(chunks)
                        all_ids.extend(ids)
                        total_chunks += len(chunks)
                        processed_files.append(doc_path)
                        logger.info("Prepared %d chunks from %s", len(chunks), doc_path)
                    else:
                        logger.warning("No chunks from %s", doc_path)

        # Single file case
        elif os.path.isfile(source):
            file_ext = Path(source).suffix.lower()
            if file_ext not in SUPPORTED_EXTENSIONS and file_ext:
                logger.warning("File extension %s not in supported list, attempting anyway: %s", file_ext, source)

            logger.info("Processing single file: %s", source)
            processed_files.append(source)
            chunks, ids = await loop.run_in_executor(None, document_to_chunks, source)
            all_chunks.extend(chunks)
            all_ids.extend(ids)
            total_chunks += len(chunks)

        else:
            return {"status": "error", "message": f"Invalid source: {source}. Must be a document file, folder, or URL."}

        if all_chunks:
            # Batch write to vectorstore in one call for performance
            logger.info("Adding %d chunks to vectorstore (batched)", len(all_chunks))
            # The vectorstore call is blocking; run in executor to avoid blocking loop
            def _add_and_persist(docs: List[Document], ids_list: List[str]) -> None:
                vectorstore.add_documents(documents=docs, ids=ids_list)
                try:
                    # persist if supported
                    if hasattr(vectorstore, "persist"):
                        vectorstore.persist()
                except Exception:
                    logger.debug("vectorstore.persist failed or not supported", exc_info=True)

            await loop.run_in_executor(None, _add_and_persist, all_chunks, all_ids)
        else:
            logger.info("No chunks generated; nothing to add to vectorstore")

        return {
            "status": "success",
            "chunks_added": total_chunks,
            "files_processed": len(processed_files),
            "files": processed_files
        }

    except Exception as e:
        logger.exception("ingest_document failed: %s", e)
        return {"status": "error", "message": str(e)}


# Legacy tool for backwards compatibility
@mcp.tool()
async def ingest_pdf(source: str) -> Dict[str, Any]:
    """
    Legacy tool. Use ingest_document instead.
    Ingest PDF(s) from URL, folder or file path.
    """
    return await ingest_document(source)


@mcp.tool()
async def retrieve(query: str, n: int = 5) -> List[Dict[str, Any]]:
    """
    Retrieve top N chunks for given query. Clamps n to reasonable bounds.
    """
    try:
        k = max(1, min(100, int(n)))
        # blocking call; run in executor
        loop = asyncio.get_running_loop()
        def _search():
            return vectorstore.similarity_search_with_score(query, k=k)
        results = await loop.run_in_executor(None, _search)

        out: List[Dict[str, Any]] = []
        for doc, score in results:
            # score is distance; convert to similarity if possible
            try:
                similarity = float(1 - score)
            except Exception:
                similarity = float(score)
            out.append({
                "text": doc.page_content,
                "metadata": doc.metadata,
                "similarity_score": similarity,
                "distance": float(score)
            })
        return out

    except Exception as e:
        logger.exception("retrieve failed: %s", e)
        return [{"error": str(e)}]


@mcp.tool()
async def db_info() -> Dict[str, Any]:
    """
    Get ChromaDB and collection information. Samples up to 100 documents for source list.
    """
    try:
        collection = getattr(vectorstore, "_collection", None)
        if collection is None:
            return {"error": "Collection object unavailable"}

        count = collection.count() if hasattr(collection, "count") else 0
        sample_size = min(100, count) if count > 0 else 0
        sources = set()

        if sample_size > 0:
            sample = collection.get(limit=sample_size, include=["metadatas"])
            metas = sample.get("metadatas") if isinstance(sample, dict) else None
            if metas:
                for metadata in metas:
                    if metadata and ("source" in metadata):
                        sources.add(metadata["source"])

        return {
            "database_path": CHROMA_PATH,
            "collection_name": COLLECTION_NAME,
            "embedding_model": EMBED_MODEL,
            "ollama_base_url": OLLAMA_BASE_URL,
            "total_chunks": count,
            "chunk_size": CHUNK_SIZE,
            "chunk_overlap": CHUNK_OVERLAP,
            "unique_sources": list(sources),
            "num_sources": len(sources)
        }

    except Exception as e:
        logger.exception("db_info failed: %s", e)
        return {"error": str(e)}


@mcp.tool()
async def check_ollama() -> Dict[str, Any]:
    """
    Check Ollama service health and embedding model availability.
    Use this tool to diagnose connection issues with Ollama.
    """
    try:
        health = check_ollama_health()

        return {
            "status": "healthy" if health["can_generate_embeddings"] else "unhealthy",
            "service_running": health["service_running"],
            "model_available": health["model_available"],
            "can_generate_embeddings": health["can_generate_embeddings"],
            "ollama_url": OLLAMA_BASE_URL,
            "embedding_model": EMBED_MODEL,
            "errors": health["errors"],
            "suggestions": health["suggestions"]
        }

    except Exception as e:
        logger.exception("check_ollama failed: %s", e)
        return {
            "status": "error",
            "error": str(e),
            "suggestions": ["Check if Ollama is installed and running"]
        }


@mcp.tool()
async def clear_db() -> Dict[str, Any]:
    """
    Clear all data from database and reinitialize an empty vectorstore.
    """
    try:
        global vectorstore
        # delete collection if supported
        try:
            vectorstore.delete_collection()
        except Exception:
            logger.debug("vectorstore.delete_collection failed or not supported", exc_info=True)

        # Remove persisted files to ensure clean reset
        try:
            if os.path.isdir(CHROMA_PATH):
                for root, dirs, files in os.walk(CHROMA_PATH, topdown=False):
                    for name in files:
                        os.remove(os.path.join(root, name))
                    for name in dirs:
                        os.rmdir(os.path.join(root, name))
        except Exception:
            logger.debug("Failed to remove chroma directory contents", exc_info=True)

        # recreate vectorstore
        vectorstore = Chroma(
            persist_directory=CHROMA_PATH,
            embedding_function=embeddings,
            collection_name=COLLECTION_NAME
        )
        try:
            if hasattr(vectorstore, "persist"):
                vectorstore.persist()
        except Exception:
            logger.debug("persist on fresh vectorstore failed or not supported", exc_info=True)

        return {"status": "success", "message": "Database cleared and reset"}

    except Exception as e:
        logger.exception("clear_db failed: %s", e)
        return {"status": "error", "message": str(e)}


# --- Main ---
def main() -> None:
    logger.info("=" * 50)
    logger.info("LangChain Chroma Vector DB Server")
    logger.info("Multi-Format Document Ingestion")
    logger.info("=" * 50)
    logger.info("ChromaDB Path: %s", CHROMA_PATH)
    logger.info("Embedding Model: %s", EMBED_MODEL)
    logger.info("Ollama URL: %s", OLLAMA_BASE_URL)
    logger.info("Chunk Size: %d", CHUNK_SIZE)
    logger.info("Chunk Overlap: %d", CHUNK_OVERLAP)
    logger.info("-" * 50)
    logger.info("Supported Formats:")
    logger.info("  PDF: %s", "✓")
    logger.info("  DOCX: %s", "✓")
    logger.info("  HTML: %s", "✓")
    logger.info("  PPTX: %s", "✓")
    logger.info("  XLSX: %s", "✓")
    logger.info("  ODT: %s", "✓")
    logger.info("  TXT/MD/CSV/JSON/XML: ✓ (built-in)")
    logger.info("-" * 50)

    # Validate Ollama connection on startup (best-effort; only meaningful if using Ollama)
    validate_ollama_connection()

    logger.info("-" * 50)
    logger.info("Available Tools:")
    logger.info("  - ingest_document: Add documents to database (all formats)")
    logger.info("  - ingest_pdf: Legacy tool for PDF ingestion")
    logger.info("  - retrieve: Search for chunks")
    logger.info("  - db_info: Get database info")
    logger.info("  - check_ollama: Check Ollama service health")
    logger.info("  - clear_db: Clear entire database")
    logger.info("=" * 50)


if __name__ == "__main__":
    main()
    mcp.run(transport="stdio")